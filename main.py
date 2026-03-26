# 백엔드 서버 실행 및 파일 요약 API를 담당하는 메인 모듈
import os
import httpx
import io
import tempfile
from pathlib import Path
from fastapi import FastAPI, File, UploadFile, HTTPException
from pydantic import BaseModel
from dotenv import load_dotenv
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
from pypdf import PdfReader
from pptx import Presentation

# from openai import OpenAI

# 로컬 개발/배포 환경에서 .env 값을 자동으로 로드
load_dotenv() # .env 파일에 있는 환경변수 불러오기

# 프로젝트 폴더 경로
BASE_DIR = Path(__file__).resolve().parent

# FastAPI 앱 인스턴스 생성
app = FastAPI()
# 정적 파일(CSS/JS) 서빙 경로
app.mount("/static", StaticFiles(directory=BASE_DIR / "static"), name="static")

# OpenAI API 사용 시 아래 두 줄 주석 해제
# from openai import OpenAI
# client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# 로컬 LLM(Ollama 등) 호출 주소와 모델명 설정
# 환경변수 미설정 시 기본값으로 동작하도록 fallback 제공
LOCAL_LLM_ENDPOINT = os.getenv("LOCAL_LLM_ENDPOINT", "http://localhost:11434/api/generate")
LOCAL_LLM_MODEL = os.getenv("LOCAL_LLM_MODEL", "qwen2.5:14b")

@app.get("/")
def root():
    # 사용자가 루트 경로로 접속하면 프론트엔드 페이지를 그대로 반환
    return FileResponse(BASE_DIR / "templates" / "index.html")

# PDF/PPTX 파일에서 텍스트 추출 함수
def extract_text_from_pdf(file_bytes: bytes) -> str:
    # 메모리 버퍼(io.BytesIO)로 PDF를 읽어 디스크 I/O를 줄임
    reader = PdfReader(io.BytesIO(file_bytes))
    pages_text = []
    # 페이지 순서대로 텍스트를 모아 하나의 문자열로 결합
    for page in reader.pages:
        # extract_text()가 None일 수 있으므로 빈 문자열로 보정
        pages_text.append(page.extract_text() or "")
    return "\n".join(pages_text).strip()


def extract_text_from_pptx(temp_path: Path) -> str:
    # python-pptx는 파일 경로 기반으로 읽는 편이 안정적이어서 임시 파일 경로 사용
    presentation = Presentation(str(temp_path))
    slide_texts = []
    # 각 슬라이드의 텍스트를 순서대로 수집
    for slide in presentation.slides:
        shape_texts = []
        for shape in slide.shapes:
            # 텍스트 속성이 있는 shape만 추출
            if hasattr(shape, "text") and shape.text:
                shape_texts.append(shape.text)
        if shape_texts:
            slide_texts.append("\n".join(shape_texts))
    return "\n\n".join(slide_texts).strip()


# 요약 요청을 처리하는 API 엔드포인트
# 프론트엔드에서 긴 텍스트를 보내면 로컬 LLM으로 요약 결과 반환
@app.post("/summarize")
async def summarize_file(file: UploadFile = File(...)):
    # 허용 MIME 타입 맵: 클라이언트가 보내는 content-type 검증에 사용
    allowed_types = {
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        "application/vnd.ms-powerpoint": "ppt",
    }
    # 파일명 확장자와 MIME 타입을 함께 확인하여 검증 정확도 향상
    extension = Path(file.filename or "").suffix.lower()
    detected_type = allowed_types.get(file.content_type)

    # 확장자/MIME 모두 허용 목록이 아니면 즉시 차단
    if extension not in {".pdf", ".pptx", ".ppt"} and detected_type not in {"pdf", "pptx", "ppt"}:
        raise HTTPException(status_code=400, detail="PDF 또는 PPT(PPTX) 파일만 업로드할 수 있습니다.")

    # 업로드 바이너리를 메모리로 읽음
    file_bytes = await file.read()
    if not file_bytes:
        raise HTTPException(status_code=400, detail="업로드된 파일이 비어 있습니다.")

    # PPTX 처리 시 생성되는 임시 파일 경로를 finally에서 정리하기 위해 보관
    temp_file_path: Path | None = None
    try:
        # 파일 유형별 텍스트 추출 분기
        if extension == ".pdf" or detected_type == "pdf":
            text = extract_text_from_pdf(file_bytes)
        elif extension == ".pptx" or detected_type == "pptx":
            # PPTX는 임시 파일로 저장 후 python-pptx로 파싱
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx", dir=BASE_DIR) as temp_file:
                temp_file.write(file_bytes)
                temp_file_path = Path(temp_file.name)
            text = extract_text_from_pptx(temp_file_path)
        else:
            # .ppt(구형 포맷) 지원 불가 안내
            raise HTTPException(
                status_code=400,
                detail="기존 PPT(.ppt) 형식은 현재 지원하지 않습니다. PPTX로 변환 후 업로드해 주세요.",
            )

        # 텍스트가 전혀 추출되지 않으면 요약을 진행하지 않음
        if not text:
            raise HTTPException(status_code=400, detail="파일에서 추출된 텍스트가 없습니다.")
    finally:
        # 예외가 발생해도 임시 파일은 반드시 삭제
        if temp_file_path and temp_file_path.exists():
            temp_file_path.unlink()
            
    # 모델에 전달할 역할/출력 형식 지시문 + 사용자 원문 결합
    prompt = (
        "너는 자료를 핵심만 요약해 주는 역할이야. 각 줄은 간결하게 작성해."
        "반드시 한국어로 대답해.\n\n"
        f"[강의자료]\n{text}"
    )

    # 로컬 LLM 호출
    try:
        # Ollama generate API 형식으로 요청
        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(
                LOCAL_LLM_ENDPOINT,
                json={
                    "model": LOCAL_LLM_MODEL,
                    "prompt": prompt,
                    "stream": False,
                },
            )
        response.raise_for_status()
        data = response.json()
        # Ollama 응답의 실제 생성 텍스트는 response 키에 담김
        summary = data.get("response", "요약 결과가 비어 있습니다.")
        # 프론트엔드가 바로 렌더링할 수 있는 JSON 형태로 반환
        return {"summary": summary}
    except Exception as exc:
        # 프론트에서 바로 처리할 수 있도록 HTTP 에러로 변환
        raise HTTPException(status_code=500, detail=f"로컬 LLM 호출 실패: {exc}")

    # OpenAI API로 다시 전환할 때 아래 블록 사용
    # response = client.chat.completions.create(
    #     model="gpt-4o-mini",
    #     messages=[
    #         {
    #             "role": "system",
    #             "content": "너는 대학생의 전공 강의자료를 핵심만 3줄로 요약해 주는 역할이야. 반드시 한국어로 대답해.",
    #         },
    #         {"role": "user", "content": request.text},
    #     ],
    # )
    # return {"summary": response.choices[0].message.content}